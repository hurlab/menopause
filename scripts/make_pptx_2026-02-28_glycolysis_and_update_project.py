#!/usr/bin/env python
"""Build/update PPTX deliverables (2026-02-28).

1) New glycolysis hypothesis deck (theme borrowed from existing project deck)
2) Updated project summary deck copied from v4.1.2 with new slides appended

Notes:
- Keep edits conservative on the existing project deck (mostly append-only).
- New slides use TEXT_TO_FIT_SHAPE where possible to avoid overflow.
"""

from __future__ import annotations

import datetime as _dt
import os
import shutil
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt
import csv


ROOT = Path(__file__).resolve().parents[1]
DOCS = ROOT / "docs"
OUT = ROOT / "docs"

TEMPLATE_PPTX = DOCS / "menopause_sns_gtex_enhanced_2026-02-02_v4.1.2.pptx"

# Figures produced in this repo (already generated).
FIG = ROOT / "GTEx_v10_AT_analysis_out" / "figs"
FIG_GLY = FIG / "next_steps_2026-02-28_glycolysis"
FIG_REC = FIG / "next_steps_2026-02-28_receptor_heatmaps"
FIG_KEN = FIG / "next_steps_2026-02-28_kenichi_panel_by_sex"
FIG_CIB = FIG / "next_steps_2026-02-28_cibersort_fractions"


def _now_stamp() -> str:
    return _dt.datetime.now().strftime("%Y-%m-%d")


def _delete_all_slides(prs: Presentation) -> None:
    # python-pptx has no public API for this; use the standard drop-rel recipe.
    # Keep removing slide 0 until empty.
    while len(prs.slides) > 0:
        sldId = prs.slides._sldIdLst[0]
        rId = sldId.rId
        prs.slides._sldIdLst.remove(sldId)
        prs.part.drop_rel(rId)


def _set_textbox_autofit(shape) -> None:
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    tf.word_wrap = True
    try:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except Exception:
        pass


def _add_bullets(slide, placeholder, bullets: list[str], font_pt: int = 20) -> None:
    tf = placeholder.text_frame
    tf.clear()
    tf.word_wrap = True
    try:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except Exception:
        pass

    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(font_pt)


def _add_picture_fit(slide, image_path: Path, left, top, width, height):
    # Fit image into the given box while preserving aspect ratio.
    im = Image.open(image_path)
    iw, ih = im.size
    im.close()

    box_w = int(width)
    box_h = int(height)

    aspect = iw / ih
    box_aspect = box_w / box_h

    if aspect >= box_aspect:
        pic_w = box_w
        pic_h = int(box_w / aspect)
    else:
        pic_h = box_h
        pic_w = int(box_h * aspect)

    pic_left = int(left + (box_w - pic_w) / 2)
    pic_top = int(top + (box_h - pic_h) / 2)

    return slide.shapes.add_picture(str(image_path), pic_left, pic_top, width=pic_w, height=pic_h)


def build_glycolysis_hypothesis_deck(out_path: Path) -> None:
    prs = Presentation(str(TEMPLATE_PPTX))
    _delete_all_slides(prs)

    # Layout indices are inherited from the template.
    L_TITLE = 0
    L_TITLE_BODY = 2
    L_TITLE_ONLY = 5

    # Slide 1: Title
    s = prs.slides.add_slide(prs.slide_layouts[L_TITLE])
    s.shapes.title.text = "Menopause, SNS remodeling, and glycolysis in human adipose"
    sub = s.placeholders[1]
    sub.text = "Hypotheses, current GTEx v10 baseline readout, and validation plan"
    date_box = None
    for sh in s.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip() == "2026-02-02":
            date_box = sh
            break
    if date_box is not None:
        date_box.text_frame.text = _now_stamp()

    # Slide 2: The issue + caution
    s = prs.slides.add_slide(prs.slide_layouts[L_TITLE_BODY])
    s.shapes.title.text = "Why glycolysis is tricky here (what we know / don’t know)"
    body = s.placeholders[1]
    _add_bullets(
        s,
        body,
        [
            "Meeting summary explicitly discussed an mRNA vs protein paradox for lipolysis (ATGL/PNPLA2, HSL/LIPE) in an HFD context; not glycolysis.",
            "GTEx adipose is baseline bulk RNA-seq: no acute stimulation, no phospho/proteomics; transcript != flux.",
            "Apparent pathway changes can be driven by cell composition (adipocyte vs SVF/immune/endothelial).",
            "So: interpret glycolysis at (1) gene nodes and (2) cell-compartment level, not only bulk gene-set scores.",
        ],
        font_pt=20,
    )

    # Slide 3: Literature context (menopause/adipose/mitochondria/glycolysis framing)
    s = prs.slides.add_slide(prs.slide_layouts[L_TITLE_BODY])
    s.shapes.title.text = "Literature context (high-level, relevant to this hypothesis)"
    body = s.placeholders[1]
    _add_bullets(
        s,
        body,
        [
            "Menopause/estrogen decline is linked to changes in fat distribution and cardiometabolic risk (review literature).",
            "Estrogen signaling is broadly connected to mitochondrial function/oxidative metabolism in metabolic tissues (review literature).",
            "Adipose aging/remodeling is commonly associated with inflammation, fibrosis, and altered metabolic programs (adipose aging reviews).",
            "Immunometabolism: inflammatory activation can increase reliance on glycolysis in immune cells (Nat Rev Immunol review literature).",
            "Adipose hypoxia/HIF-1alpha programs can promote glycolysis and fibrosis in obesity/remodeling contexts (review literature).",
        ],
        font_pt=20,
    )

    # Slide 4: Literature context (SNS/adrenergic/catecholamine resistance)
    s = prs.slides.add_slide(prs.slide_layouts[L_TITLE_BODY])
    s.shapes.title.text = "Literature context (SNS/adrenergic signaling and interpretation)"
    body = s.placeholders[1]
    _add_bullets(
        s,
        body,
        [
            "SNS drives lipolysis via beta-adrenergic receptors → cAMP/PKA → phosphorylation-driven activation (minutes), not transcription alone.",
            "Chronic catecholamine exposure can cause adrenergic receptor desensitization/downregulation (catecholamine resistance concept).",
            "Therefore baseline bulk mRNA can reflect chronic adaptation and/or composition shifts rather than acute functional flux.",
            "Our GTEx readout is baseline; acute immediate-early gene signatures are expected to be weak at baseline.",
        ],
        font_pt=20,
    )

    # Slide 5: Working framing schematic (simple box/arrow diagram)
    s = prs.slides.add_slide(prs.slide_layouts[L_TITLE_ONLY])
    s.shapes.title.text = "Working hypothesis: menopause → adrenergic remodeling + remodeling of tissue composition"
    title = s.shapes.title
    top0 = title.top + title.height + Inches(0.2)
    left0 = Inches(0.6)
    w = prs.slide_width - Inches(1.2)
    h = prs.slide_height - top0 - Inches(0.4)

    # Diagram coordinates (within the box). Keep within 10x5.625 slide bounds (template).
    def box(x, y, bw, bh, text, fill_rgb=(240, 248, 255)):
        shp = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, bw, bh)
        shp.fill.solid()
        shp.fill.fore_color.rgb = RGBColor(*fill_rgb)
        shp.line.color.rgb = RGBColor(60, 60, 60)
        tf = shp.text_frame
        tf.text = text
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.word_wrap = True
        try:
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        except Exception:
            pass
        return shp

    def arrow(x1, y1, x2, y2):
        ln = s.shapes.add_connector(1, x1, y1, x2, y2)
        ln.line.width = Pt(2)
        return ln

    # Vertical flow layout (robust to template slide size).
    row_h = Inches(0.85)
    gap = Inches(0.18)
    x = left0
    bw = w

    y1 = top0
    b1 = box(x, y1, bw, row_h, "Menopause / estrogen decline (+ aging)")
    y2 = y1 + row_h + gap
    b2 = box(x, y2, bw, row_h, "Chronic SNS tone changes + adrenergic remodeling\n(catecholamine resistance / desensitization modules)")
    y3 = y2 + row_h + gap
    b3 = box(x, y3, bw, row_h, "Tissue remodeling in adipose\n(cell composition shifts in bulk: adipocyte ↓; SVF/immune/endothelial/fibro ↑)", fill_rgb=(255, 250, 240))
    y4 = y3 + row_h + gap
    b4 = box(x, y4, bw, row_h, "Observed bulk transcript patterns (GTEx baseline)\nOXPHOS ↓; inflammatory response ↑; adipogenesis ↓; glycolysis mixed/flat overall", fill_rgb=(245, 255, 250))

    arrow(b1.left + b1.width / 2, b1.top + b1.height, b2.left + b2.width / 2, b2.top)
    arrow(b2.left + b2.width / 2, b2.top + b2.height, b3.left + b3.width / 2, b3.top)
    arrow(b3.left + b3.width / 2, b3.top + b3.height, b4.left + b4.width / 2, b4.top)

    # Slide 6: Data/methods
    s = prs.slides.add_slide(prs.slide_layouts[L_TITLE_BODY])
    s.shapes.title.text = "Data + scoring used for this quick glycolysis check"
    body = s.placeholders[1]
    _add_bullets(
        s,
        body,
        [
            "GTEx v10 adipose bulk RNA-seq (subcutaneous vs visceral), baseline tissue.",
            "Scores computed like the rest of this repo: VST per depot; per-gene z within (depot+sex); gene-set score = mean(z).",
            "Female menopause-proxy bins (age midpoint): pre <45, peri 45–55, post >55.",
            "Linear models adjusted for sequencing center (SMCENTER); sensitivity with CIBERSORT fraction covariates.",
        ],
        font_pt=20,
    )

    # Slide 7: Glycolysis hallmark boxplots (subcutaneous + visceral)
    s = prs.slides.add_slide(prs.slide_layouts[L_TITLE_ONLY])
    s.shapes.title.text = "Result: Hallmark Glycolysis is not robustly shifted (bulk, baseline)"
    title = s.shapes.title
    top = title.top + title.height + Inches(0.2)
    margin = Inches(0.5)
    avail_w = prs.slide_width - margin * 2
    avail_h = prs.slide_height - top - Inches(0.4)

    # two panels side-by-side
    img1 = FIG_GLY / "female_menopause_groups_hallmarks_boxplots.png"
    img2 = FIG_GLY / "female_menopause_groups_hallmarks_boxplots_visceral.png"
    box_w = int(avail_w / 2 - Inches(0.2))
    _add_picture_fit(s, img1, margin, top, box_w, avail_h)
    _add_picture_fit(s, img2, margin + box_w + Inches(0.4), top, box_w, avail_h)

    # Slide 8: Model summary table (key stats)
    s = prs.slides.add_slide(prs.slide_layouts[L_TITLE_ONLY])
    s.shapes.title.text = "Model summary (female pre/peri/post; SMCENTER-adjusted; effect sizes are z-score units)"
    title = s.shapes.title
    top = title.top + title.height + Inches(0.15)
    margin = Inches(0.55)

    stats_path = (
        ROOT
        / "GTEx_v10_AT_analysis_out"
        / "tables"
        / "next_steps_2026-02-28_glycolysis"
        / "female_pre_peri_post_lm_stats_hallmarks.tsv"
    )
    stats_vis_path = (
        ROOT
        / "GTEx_v10_AT_analysis_out"
        / "tables"
        / "next_steps_2026-02-28_glycolysis"
        / "female_pre_peri_post_lm_stats_hallmarks_visceral.tsv"
    )

    def load_rows(path: Path, depot: str, gene_sets: list[str], model: str):
        out = {}
        with path.open() as f:
            rdr = csv.DictReader(f, delimiter="\t")
            for r in rdr:
                if r["depot"] != depot:
                    continue
                if r["model"] != model:
                    continue
                gs = r["gene_set"]
                if gs not in gene_sets:
                    continue
                out[gs] = r
        return out

    key_sets = [
        "HALLMARK_GLYCOLYSIS",
        "HALLMARK_OXIDATIVE_PHOSPHORYLATION",
        "HALLMARK_INFLAMMATORY_RESPONSE",
        "HALLMARK_ADIPOGENESIS",
    ]

    sub_base = load_rows(stats_path, "subcutaneous", key_sets, "base")
    sub_adj = load_rows(stats_path, "subcutaneous", key_sets, "adj_cibersort")
    vis_base = load_rows(stats_vis_path, "visceral", key_sets, "base")
    vis_adj = load_rows(stats_vis_path, "visceral", key_sets, "adj_cibersort")

    # Create a compact table: rows = key gene sets; cols = sub/vis x base/adj (post vs pre est,p)
    nrows = 1 + len(key_sets)
    ncols = 5
    tbl = s.shapes.add_table(nrows, ncols, margin, top, prs.slide_width - margin * 2, Inches(3.2)).table
    hdr = [
        "Gene set",
        "Subq base\npost-pre",
        "Subq +CIBERSORT\npost-pre",
        "Visc base\npost-pre",
        "Visc +CIBERSORT\npost-pre",
    ]
    for j, htxt in enumerate(hdr):
        c = tbl.cell(0, j)
        c.text = htxt
        c.text_frame.paragraphs[0].font.size = Pt(12)
        c.text_frame.paragraphs[0].font.bold = True

    def fmt(r):
        if not r:
            return "NA"
        est = float(r["post_vs_pre_est"])
        p = r.get("post_vs_pre_p_fmt") or r.get("post_vs_pre_p")
        return f"{est:+.3f}; p={p}"

    for i, gs in enumerate(key_sets, start=1):
        tbl.cell(i, 0).text = gs.replace("HALLMARK_", "")
        tbl.cell(i, 1).text = fmt(sub_base.get(gs))
        tbl.cell(i, 2).text = fmt(sub_adj.get(gs))
        tbl.cell(i, 3).text = fmt(vis_base.get(gs))
        tbl.cell(i, 4).text = fmt(vis_adj.get(gs))
        for j in range(ncols):
            tf = tbl.cell(i, j).text_frame
            for p in tf.paragraphs:
                p.font.size = Pt(12)

    # Slide note under table
    note_h = Inches(0.55)
    note_top = prs.slide_height - Inches(0.25) - note_h
    tx = s.shapes.add_textbox(margin, note_top, prs.slide_width - margin * 2, note_h)
    tf = tx.text_frame
    tf.text = "Key readout: Glycolysis is not significant in these models; OXPHOS/adipogenesis/inflammation show clearer subcutaneous shifts."
    tf.paragraphs[0].font.size = Pt(14)
    _set_textbox_autofit(tx)

    # Slide 9: Trajectories by sex (context)
    s = prs.slides.add_slide(prs.slide_layouts[L_TITLE_ONLY])
    s.shapes.title.text = "Context: trajectories by sex across age bins (hallmarks panel)"
    title = s.shapes.title
    top = title.top + title.height + Inches(0.2)
    margin = Inches(0.5)
    avail_w = prs.slide_width - margin * 2
    avail_h = prs.slide_height - top - Inches(0.4)

    img1 = FIG_GLY / "sex_agebin_trajectories_hallmarks.png"
    img2 = FIG_GLY / "sex_agebin_trajectories_hallmarks_visceral.png"
    box_w = int(avail_w / 2 - Inches(0.2))
    _add_picture_fit(s, img1, margin, top, box_w, avail_h)
    _add_picture_fit(s, img2, margin + box_w + Inches(0.4), top, box_w, avail_h)

    # Slide 10: What *did* move (OXPHOS/inflammation/adipogenesis)
    s = prs.slides.add_slide(prs.slide_layouts[L_TITLE_BODY])
    s.shapes.title.text = "Interpretation: glycolysis looks flat, but other programs shift"
    body = s.placeholders[1]
    _add_bullets(
        s,
        body,
        [
            "In the same run (female, subcutaneous): OXPHOS down (post vs pre), adipogenesis down, inflammatory response up.",
            "These effects are partly sensitive to composition adjustment, consistent with remodeling in bulk.",
            "A ‘glycolysis discrepancy’ may therefore be gene-node-specific and/or cell-compartment-specific, not a large bulk hallmark shift.",
        ],
        font_pt=22,
    )

    # Slide 11: Composition evidence (CIBERSORT)
    s = prs.slides.add_slide(prs.slide_layouts[L_TITLE_ONLY])
    s.shapes.title.text = "Cell composition (first pass): CIBERSORT fractions vary with age/sex"
    title = s.shapes.title
    top = title.top + title.height + Inches(0.2)
    margin = Inches(0.5)
    avail_w = prs.slide_width - margin * 2
    avail_h = prs.slide_height - top - Inches(0.4)

    img1 = FIG_CIB / "cibersort_subcutaneous_perm0_topN50_minCells200__fractions_mean_by_sex_agebin.png"
    img2 = FIG_CIB / "cibersort_visceral_perm0_topN50_minCells200__fractions_mean_by_sex_agebin.png"
    box_w = int(avail_w / 2 - Inches(0.2))
    _add_picture_fit(s, img1, margin, top, box_w, avail_h)
    _add_picture_fit(s, img2, margin + box_w + Inches(0.4), top, box_w, avail_h)

    # Slide 12: Hypotheses
    s = prs.slides.add_slide(prs.slide_layouts[L_TITLE_BODY])
    s.shapes.title.text = "Hypotheses to reconcile ‘glycolysis discrepancies’"
    body = s.placeholders[1]
    _add_bullets(
        s,
        body,
        [
            "H1 (mixture): gene-level glycolysis shifts reflect SVF/immune/endothelial composition changes more than adipocyte intrinsic rewiring.",
            "H2 (compartment): adipocyte glycolysis may change at key nodes (e.g., GLUT4/SLC2A4) without a strong bulk hallmark shift.",
            "H3 (remodeling): menopause-linked adrenergic remodeling/catecholamine resistance alters substrate use and mitochondrial programs (OXPHOS down), with secondary glycolysis effects.",
        ],
        font_pt=20,
    )

    # Slide 13: Concrete predictions
    s = prs.slides.add_slide(prs.slide_layouts[L_TITLE_BODY])
    s.shapes.title.text = "Concrete predictions (computational + experimental)"
    body = s.placeholders[1]
    _add_bullets(
        s,
        body,
        [
            "If H1: glycolysis score correlates with macrophage/endothelial/ASPC fractions; group effects shrink after fraction adjustment.",
            "If H2: adipocyte-enriched nodes (SLC2A4, IRS1/AKT axis) show clearer group effects than generic glycolysis enzymes.",
            "If H3: OXPHOS down + inflammation up is reproducible; adrenergic responsiveness assays show desensitization (cAMP/glycerol release).",
        ],
        font_pt=20,
    )

    # Slide 14: Next computational analyses
    s = prs.slides.add_slide(prs.slide_layouts[L_TITLE_BODY])
    s.shapes.title.text = "Next computational work (low lift, highest ROI)"
    body = s.placeholders[1]
    _add_bullets(
        s,
        body,
        [
            "Gene-node analysis: SLC2A4/SLC2A1, HK2, PFKM/PFKP, PKM, LDHA, SLC16A1/SLC16A3 by depot/sex/age.",
            "Modeling: repeat ~ SMCENTER + group, then + CIBERSORT subset; test glycolysis-vs-fraction correlations.",
            "If needed: move from ‘fractions as covariates’ to higher-resolution deconvolution (CIBERSORTx HiRes) with a stronger adipose atlas reference.",
        ],
        font_pt=20,
    )

    # Slide 15: Minimal experimental validation
    s = prs.slides.add_slide(prs.slide_layouts[L_TITLE_BODY])
    s.shapes.title.text = "Fast experimental validations (one clean figure)"
    body = s.placeholders[1]
    _add_bullets(
        s,
        body,
        [
            "Separate adipocytes vs SVF when possible (or quantify composition markers).",
            "Assays: glucose uptake (basal + insulin), lactate secretion, Seahorse ECAR/OCR.",
            "Adrenergic responsiveness: isoproterenol-stimulated cAMP + glycerol release; ADRB2/ADRB3 protein if feasible.",
            "qPCR/Western panel: SLC2A4, HK2, LDHA + inflammation markers + mitochondrial/OXPHOS markers.",
        ],
        font_pt=20,
    )

    # Slide 16: Selected literature (curated)
    s = prs.slides.add_slide(prs.slide_layouts[L_TITLE_BODY])
    s.shapes.title.text = "Selected literature (for framing; not exhaustive)"
    body = s.placeholders[1]
    refs = [
        "Karpe F, Pinnick KE. Nat Rev Endocrinol (2015) upper vs lower body adipose review.",
        "Mauvais-Jarvis F. Biol Sex Differ (2015) sex differences in metabolic homeostasis review.",
        "Palmer AK, Kirkland JL. Exp Gerontol (2016) adipose aging review.",
        "O’Neill LAJ et al. Nat Rev Immunol (2016) immunometabolism review (glycolysis in inflammatory activation).",
        "Adipose hypoxia / HIF-1alpha → inflammation/fibrosis/glycolysis programs (review literature).",
        "Adrenergic receptor desensitization/catecholamine resistance concepts in adipose (review literature).",
        "Project-specific: GTEx v10 baseline bulk; fraction adjustment/deconvolution needed for glycolysis interpretation.",
    ]
    _add_bullets(s, body, refs, font_pt=16)

    # Slide 17: Appendix genes
    s = prs.slides.add_slide(prs.slide_layouts[L_TITLE_BODY])
    s.shapes.title.text = "Appendix: glycolysis/transport nodes to prioritize"
    body = s.placeholders[1]
    _add_bullets(
        s,
        body,
        [
            "Transport: SLC2A4 (GLUT4), SLC2A1 (GLUT1)",
            "Core enzymes: HK2, PFKM/PFKP, ALDOA, GAPDH, ENO1, PKM, LDHA",
            "Lactate transport: SLC16A1 (MCT1), SLC16A3 (MCT4)",
            "Insulin axis context: IRS1, AKT2 (if needed)",
        ],
        font_pt=22,
    )

    prs.save(str(out_path))


def update_project_summary_deck(out_path: Path) -> None:
    src = TEMPLATE_PPTX
    tmp = out_path
    shutil.copy2(src, tmp)

    prs = Presentation(str(tmp))

    # Update title slide date text.
    s1 = prs.slides[0]
    for sh in s1.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip() == "2026-02-02":
            sh.text_frame.text = _now_stamp()

    # Fix a known textbox overflow already present in v4.1.2 (slide 21). Keep look; just keep it in-bounds.
    try:
        s21 = prs.slides[20]  # 0-based
        for sh in s21.shapes:
            if sh.shape_type == 17:  # TEXT_BOX
                # If it extends past bottom, shrink height to fit and enable autofit.
                if sh.top + sh.height > prs.slide_height:
                    sh.height = prs.slide_height - sh.top - Inches(0.2)
                    _set_textbox_autofit(sh)
    except Exception:
        pass

    L_TITLE_ONLY = 5
    L_TITLE_BODY = 2

    def add_two_images_slide(title: str, img_left: Path, img_right: Path, subtitle: str | None = None):
        s = prs.slides.add_slide(prs.slide_layouts[L_TITLE_ONLY])
        s.shapes.title.text = title
        title_sh = s.shapes.title
        top = title_sh.top + title_sh.height + Inches(0.15)
        margin = Inches(0.5)
        avail_w = prs.slide_width - margin * 2
        avail_h = prs.slide_height - top - Inches(0.35)
        box_w = int(avail_w / 2 - Inches(0.15))
        _add_picture_fit(s, img_left, margin, top, box_w, avail_h)
        _add_picture_fit(s, img_right, margin + box_w + Inches(0.3), top, box_w, avail_h)
        if subtitle:
            # Small note at bottom
            tx = s.shapes.add_textbox(margin, prs.slide_height - Inches(0.55), prs.slide_width - margin * 2, Inches(0.35))
            tf = tx.text_frame
            tf.text = subtitle
            tf.paragraphs[0].font.size = Pt(12)
            tf.paragraphs[0].alignment = PP_ALIGN.LEFT
            _set_textbox_autofit(tx)

    # New slide: adrenergic receptors heatmaps (beta + alpha2)
    add_two_images_slide(
        "Update (2026-02-28): adrenergic receptor expression (depot/sex/age bins)",
        FIG_REC / "heatmap_ADRB1_ADRB2_ADRB3_meanVST_by_depot_sex_agebin.png",
        FIG_REC / "heatmap_ADRA2A_ADRA2B_ADRA2C_meanVST_by_depot_sex_agebin.png",
        subtitle="Heatmaps show mean VST per (depot × sex × age bin).",
    )

    # New slide: Kenichi panel by sex (subq + visceral)
    add_two_images_slide(
        "Update (2026-02-28): Kenichi panel (Slides 7–8 style), faceted by sex",
        FIG_KEN / "figure_candidate__lipolysis_thermogenesis__by_sex.png",
        FIG_KEN / "figure_candidate__lipolysis_thermogenesis__by_sex_visceral.png",
        subtitle="Same scoring as Slides 7–8; plotted separately for male vs female.",
    )

    # New slide: CIBERSORT fractions
    add_two_images_slide(
        "Update (2026-02-28): first-pass cell fractions (CIBERSORT; GSE176171 reference)",
        FIG_CIB / "cibersort_subcutaneous_perm0_topN50_minCells200__fractions_mean_by_sex_agebin.png",
        FIG_CIB / "cibersort_visceral_perm0_topN50_minCells200__fractions_mean_by_sex_agebin.png",
        subtitle="Fractions are mean by (sex × age bin); perm=0 first pass; TopN50 signature.",
    )

    # New slide: glycolysis / hallmarks
    add_two_images_slide(
        "Update (2026-02-28): glycolysis and related hallmarks (baseline GTEx)",
        FIG_GLY / "female_menopause_groups_hallmarks_boxplots.png",
        FIG_GLY / "sex_agebin_trajectories_hallmarks.png",
        subtitle="Key point: Hallmark Glycolysis shows no robust pre/peri/post shift; OXPHOS/inflammation/adipogenesis move.",
    )

    # New slide: updated file list
    s = prs.slides.add_slide(prs.slide_layouts[L_TITLE_BODY])
    s.shapes.title.text = "Update (2026-02-28): new outputs added since v4.1.2"
    body = s.placeholders[1]
    _add_bullets(
        s,
        body,
        [
            "Adrenergic receptor heatmaps: figs/next_steps_2026-02-28_receptor_heatmaps/*.png",
            "Sex-faceted Kenichi panels: figs/next_steps_2026-02-28_kenichi_panel_by_sex/*.png",
            "CIBERSORT first-pass fractions: figs+tables/next_steps_2026-02-28_cibersort_fractions/*",
            "Glycolysis/Hypoxia/OXPHOS hallmarks: figs+tables/next_steps_2026-02-28_glycolysis/*",
            "Memo: docs/glycolysis_hypothesis_memo_2026-02-28.md",
        ],
        font_pt=18,
    )

    prs.save(str(tmp))


def update_project_summary_deck_v4_2(out_path: Path) -> None:
    # Build from v4.1.2 to preserve all existing formatting and slides 1-23.
    src = TEMPLATE_PPTX
    shutil.copy2(src, out_path)

    prs = Presentation(str(out_path))

    # Update title slide date text.
    s1 = prs.slides[0]
    for sh in s1.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip() == "2026-02-02":
            sh.text_frame.text = _now_stamp()

    # Fix a known textbox overflow already present in the base deck (slide 21).
    try:
        s21 = prs.slides[20]  # 0-based
        for sh in s21.shapes:
            if sh.shape_type == 17:  # TEXT_BOX
                if sh.top + sh.height > prs.slide_height:
                    sh.height = prs.slide_height - sh.top - Inches(0.2)
                    _set_textbox_autofit(sh)
    except Exception:
        pass

    L_TITLE_ONLY = 5
    L_TITLE_BODY = 2

    def add_two_images_slide(title: str, img_left: Path, img_right: Path, subtitle: str | None = None):
        s = prs.slides.add_slide(prs.slide_layouts[L_TITLE_ONLY])
        s.shapes.title.text = title
        title_sh = s.shapes.title
        top = title_sh.top + title_sh.height + Inches(0.15)
        margin = Inches(0.5)
        avail_w = prs.slide_width - margin * 2
        avail_h = prs.slide_height - top - Inches(0.35)
        box_w = int(avail_w / 2 - Inches(0.15))
        _add_picture_fit(s, img_left, margin, top, box_w, avail_h)
        _add_picture_fit(s, img_right, margin + box_w + Inches(0.3), top, box_w, avail_h)
        if subtitle:
            tx = s.shapes.add_textbox(margin, prs.slide_height - Inches(0.55), prs.slide_width - margin * 2, Inches(0.35))
            tf = tx.text_frame
            tf.text = subtitle
            tf.paragraphs[0].font.size = Pt(12)
            tf.paragraphs[0].alignment = PP_ALIGN.LEFT
            _set_textbox_autofit(tx)

    def add_one_image_slide(title: str, img: Path, subtitle: str | None = None):
        s = prs.slides.add_slide(prs.slide_layouts[L_TITLE_ONLY])
        s.shapes.title.text = title
        title_sh = s.shapes.title
        top = title_sh.top + title_sh.height + Inches(0.15)
        margin = Inches(0.5)
        avail_w = prs.slide_width - margin * 2
        avail_h = prs.slide_height - top - Inches(0.35)
        _add_picture_fit(s, img, margin, top, int(avail_w), int(avail_h))
        if subtitle:
            tx = s.shapes.add_textbox(margin, prs.slide_height - Inches(0.55), prs.slide_width - margin * 2, Inches(0.35))
            tf = tx.text_frame
            tf.text = subtitle
            tf.paragraphs[0].font.size = Pt(12)
            tf.paragraphs[0].alignment = PP_ALIGN.LEFT
            _set_textbox_autofit(tx)

    # Update block starts here (replaces the v4.1.3 appended set with expanded v4.2 set).

    # 1) Receptor heatmaps (group means; as before)
    add_two_images_slide(
        "Update (2026-02-28): adrenergic receptor expression (depot/sex/age bins)",
        FIG_REC / "heatmap_ADRB1_ADRB2_ADRB3_meanVST_by_depot_sex_agebin.png",
        FIG_REC / "heatmap_ADRA2A_ADRA2B_ADRA2C_meanVST_by_depot_sex_agebin.png",
        subtitle="Heatmaps show mean VST per (depot × sex × age bin).",
    )

    # 1b) Receptor heatmaps with values in cells
    add_two_images_slide(
        "Update (2026-02-28): adrenergic receptor heatmaps with mean VST values (2 decimals)",
        FIG_REC / "heatmap_ADRB1_ADRB2_ADRB3_meanVST_by_depot_sex_agebin_with_values.png",
        FIG_REC / "heatmap_ADRA2A_ADRA2B_ADRA2C_meanVST_by_depot_sex_agebin_with_values.png",
        subtitle="Same as prior slide, but cells display mean VST values to ease cross-cell comparisons.",
    )

    # 1c) Receptor heatmaps per individual samples (ordered by age bin)
    add_two_images_slide(
        "Update (2026-02-28): adrenergic receptor expression per sample (ordered by age bin)",
        FIG_REC / "heatmap_ADRB1_ADRB2_ADRB3_VST_per_sample_ordered_by_agebin.png",
        FIG_REC / "heatmap_ADRA2A_ADRA2B_ADRA2C_VST_per_sample_ordered_by_agebin.png",
        subtitle="Per-sample VST heatmaps; samples ordered by GTEx age bin within each (sex × depot) facet.",
    )

    # 2) Kenichi panels
    add_two_images_slide(
        "Update (2026-02-28): Kenichi panel (Slides 7–8 style), faceted by sex",
        FIG_KEN / "figure_candidate__lipolysis_thermogenesis__by_sex.png",
        FIG_KEN / "figure_candidate__lipolysis_thermogenesis__by_sex_visceral.png",
        subtitle="Same scoring as Slides 7–8; plotted separately for male vs female.",
    )

    # 2b) Kenichi panels with F/M interleaved groups in one plot
    add_two_images_slide(
        "Update (2026-02-28): Kenichi panel with female/male interleaved groups (single plot)",
        FIG_KEN / "figure_candidate__lipolysis_thermogenesis__F_M_interleaved.png",
        FIG_KEN / "figure_candidate__lipolysis_thermogenesis__F_M_interleaved_visceral.png",
        subtitle="Group order: F-pre, M-pre, F-peri, M-peri, F-post, M-post (easier cross-sex comparison).",
    )

    # 2c) Interleaved Kenichi panels with medians labeled
    add_two_images_slide(
        "Update (2026-02-28): Kenichi panel (interleaved) with medians labeled",
        FIG_KEN / "figure_candidate__lipolysis_thermogenesis__F_M_interleaved_with_medians.png",
        FIG_KEN / "figure_candidate__lipolysis_thermogenesis__F_M_interleaved_with_medians_visceral.png",
        subtitle="Median values (2 decimals) are annotated above each group.",
    )

    # 3) CIBERSORT reference slide (GSE176171)
    s = prs.slides.add_slide(prs.slide_layouts[L_TITLE_ONLY])
    s.shapes.title.text = "Update (2026-02-28): CIBERSORT reference (GSE176171) used for first-pass fractions"
    title_sh = s.shapes.title
    top = title_sh.top + title_sh.height + Inches(0.15)
    margin = Inches(0.5)
    avail_h = prs.slide_height - top - Inches(0.35)
    left_w = Inches(5.4)
    right_w = prs.slide_width - margin * 2 - left_w - Inches(0.3)

    ref_img = FIG_CIB / "gse176171_reference_celltype_counts.png"
    _add_picture_fit(s, ref_img, margin, top, int(left_w), avail_h)

    # Pull a few concrete numbers from the derived TSV (if present) to reduce hand-wavy text.
    ref_counts_tsv = ROOT / "external_data" / "GSE176171" / "derived" / "reference_cell_type_counts.tsv"
    total_cells = None
    n_types = None
    top_types = []
    if ref_counts_tsv.exists():
        import csv as _csv
        with ref_counts_tsv.open() as f:
            rdr = _csv.DictReader(f, delimiter="\t")
            rows = [(r["cell_type"], int(float(r["n_cells"]))) for r in rdr]
        rows.sort(key=lambda x: x[1], reverse=True)
        total_cells = sum(n for _, n in rows)
        n_types = len(rows)
        top_types = rows[:5]

    bullets = [
        "Reference: GSE176171 (Hs10X MatrixMarket + cell-level metadata).",
        "Human cells only; broad adipose cell types (adipocyte, ASPC, endothelial, immune, etc.).",
        "Publication: Emont MP, Jacobs C, Essene AL, et al. A single-cell atlas of human and mouse white adipose tissue. Nature. 2022;603:926-933. doi:10.1038/s41586-022-04518-2.",
        "Signature: pseudo-bulk CPM per cell type; Top-N markers per type by margin; Ensembl IDs w/o version.",
        "CIBERSORT settings (first pass): TopN=50; minCells/type=200; perm=0; QN=false.",
    ]
    if total_cells is not None and n_types is not None:
        bullets.insert(
            2,
            f"Kept reference size: {total_cells:,} cells across {n_types} cell types (top types: "
            + ", ".join([f"{ct}={n:,}" for ct, n in top_types])
            + ").",
        )

    tx = s.shapes.add_textbox(margin + left_w + Inches(0.3), top, right_w, avail_h)
    tf = tx.text_frame
    tf.clear()
    tf.word_wrap = True
    try:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except Exception:
        pass
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(18)

    # 3b) Cleaned CIBERSORT fraction plots (no P-value etc) split into two slides
    add_one_image_slide(
        "Update (2026-02-28): CIBERSORT fractions (subcutaneous; cleaned to cell types only)",
        FIG_CIB / "cibersort_subcutaneous_perm0_topN50_minCells200__fractions_mean_by_sex_agebin_clean.png",
        subtitle="Cleaned: only cell-type fraction columns (excludes P-value/Correlation/RMSE).",
    )
    add_one_image_slide(
        "Update (2026-02-28): CIBERSORT fractions (visceral; cleaned to cell types only)",
        FIG_CIB / "cibersort_visceral_perm0_topN50_minCells200__fractions_mean_by_sex_agebin_clean.png",
        subtitle="Cleaned: only cell-type fraction columns (excludes P-value/Correlation/RMSE).",
    )

    # 4) Glycolysis/hallmarks slide (as before)
    add_two_images_slide(
        "Update (2026-02-28): glycolysis and related hallmarks (baseline GTEx)",
        FIG_GLY / "female_menopause_groups_hallmarks_boxplots.png",
        FIG_GLY / "sex_agebin_trajectories_hallmarks.png",
        subtitle="Key point: Hallmark Glycolysis shows no robust pre/peri/post shift; OXPHOS/inflammation/adipogenesis move.",
    )

    # 5) Updated file list slide
    s = prs.slides.add_slide(prs.slide_layouts[L_TITLE_BODY])
    s.shapes.title.text = "Update (2026-02-28): new outputs added since v4.1.2"
    body = s.placeholders[1]
    _add_bullets(
        s,
        body,
        [
            "Adrenergic receptor heatmaps: figs/next_steps_2026-02-28_receptor_heatmaps/*.png",
            "  - with values: *meanVST*_with_values.png",
            "  - per-sample: *VST_per_sample_ordered_by_agebin.png",
            "Sex-faceted Kenichi panels: figs/next_steps_2026-02-28_kenichi_panel_by_sex/*.png",
            "  - interleaved F/M: *F_M_interleaved*.png",
            "  - interleaved + medians: *F_M_interleaved_with_medians*.png",
            "CIBERSORT first-pass fractions: figs+tables/next_steps_2026-02-28_cibersort_fractions/*",
            "  - cleaned plots: *fractions_mean_by_sex_agebin_clean.png",
            "  - reference summary: gse176171_reference_celltype_counts.png",
            "Glycolysis/Hypoxia/OXPHOS hallmarks: figs+tables/next_steps_2026-02-28_glycolysis/*",
            "Memo: docs/glycolysis_hypothesis_memo_2026-02-28.md",
        ],
        font_pt=16,
    )

    prs.save(str(out_path))


def main() -> None:
    out_gly = OUT / f"menopause_glycolysis_hypothesis_{_now_stamp()}.pptx"
    out_proj = OUT / f"menopause_sns_gtex_enhanced_{_now_stamp()}_v4.1.3.pptx"
    out_proj_v42 = OUT / f"menopause_sns_gtex_enhanced_{_now_stamp()}_v4.2.pptx"
    out_proj_v421 = OUT / f"menopause_sns_gtex_enhanced_{_now_stamp()}_v4.2.1.pptx"

    # Sanity checks
    for p in [
        TEMPLATE_PPTX,
        FIG_GLY / "female_menopause_groups_hallmarks_boxplots.png",
        FIG_GLY / "female_menopause_groups_hallmarks_boxplots_visceral.png",
        FIG_GLY / "sex_agebin_trajectories_hallmarks.png",
        FIG_GLY / "sex_agebin_trajectories_hallmarks_visceral.png",
        FIG_REC / "heatmap_ADRB1_ADRB2_ADRB3_meanVST_by_depot_sex_agebin.png",
        FIG_REC / "heatmap_ADRA2A_ADRA2B_ADRA2C_meanVST_by_depot_sex_agebin.png",
        FIG_KEN / "figure_candidate__lipolysis_thermogenesis__by_sex.png",
        FIG_KEN / "figure_candidate__lipolysis_thermogenesis__by_sex_visceral.png",
        FIG_CIB / "cibersort_subcutaneous_perm0_topN50_minCells200__fractions_mean_by_sex_agebin.png",
        FIG_CIB / "cibersort_visceral_perm0_topN50_minCells200__fractions_mean_by_sex_agebin.png",
    ]:
        if not p.exists():
            raise FileNotFoundError(str(p))

    build_glycolysis_hypothesis_deck(out_gly)
    update_project_summary_deck(out_proj)
    update_project_summary_deck_v4_2(out_proj_v42)
    update_project_summary_deck_v4_2(out_proj_v421)

    print(f"Wrote: {out_gly}")
    print(f"Wrote: {out_proj}")
    print(f"Wrote: {out_proj_v42}")
    print(f"Wrote: {out_proj_v421}")


if __name__ == "__main__":
    main()
